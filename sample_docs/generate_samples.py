"""
테스트/데모용 샘플 문서를 생성하고 DB에 등록하는 스크립트.
실제 서비스에는 필요 없으며, 개발 중 동작 확인용입니다.
(예시 문서 내용은 실제 식약처 고시가 아닌 데모를 위한 가상의 예시입니다.)
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from pptx import Presentation
from pptx.util import Inches, Pt

from modules import db, extract

SAMPLE_DIR = os.path.dirname(os.path.abspath(__file__))

pdfmetrics.registerFont(UnicodeCIDFont("HYSMyeongJo-Medium"))
KOREAN_FONT = "HYSMyeongJo-Medium"


def make_notice_pdf():
    path = os.path.join(SAMPLE_DIR, "notice_sample.pdf")
    c = canvas.Canvas(path, pagesize=A4)
    lines = [
        "식품의약품안전처 고시 제2026-45호",
        "완제의약품 미생물한도 시험법 개정 고시 (예시 문서)",
        "",
        "1. 개정 배경",
        "대한민국약전 제13개정을 반영하여 완제의약품 미생물한도 시험법 중",
        "총호기성생균수 시험 배지 조성 및 판정 기준을 개정함.",
        "",
        "2. 주요 개정 내용",
        "가. 총호기성생균수 판정 기준을 100 CFU/g에서 10 CFU/g로 강화",
        "나. 시험용 배지는 카제인대두소화한천배지(TSA)로 통일",
        "다. 대장균, 살모넬라균 등 특정미생물 시험법 절차 명확화",
        "",
        "3. 시행일 및 경과조치",
        "이 고시는 2026년 10월 1일부터 시행한다.",
        "다만, 이 고시 시행 당시 이미 허가(신고)된 품목의 경우",
        "2027년 3월 31일까지 유예기간을 부여하며, 그 기간 동안은",
        "종전 기준과 개정 기준 중 하나를 선택하여 적용할 수 있다.",
        "",
        "4. 적용 범위",
        "신규 허가 신청 품목은 시행일부터 즉시 개정 기준을 적용한다.",
    ]
    text_obj = c.beginText(50, 800)
    text_obj.setFont(KOREAN_FONT, 10)
    for line in lines:
        text_obj.textLine(line)
    c.drawText(text_obj)
    c.save()
    return path


def make_qna_pdf():
    path = os.path.join(SAMPLE_DIR, "qna_sample.pdf")
    c = canvas.Canvas(path, pagesize=A4)
    lines = [
        "미생물한도 시험법 개정 관련 Q&A (예시 문서)",
        "",
        "Q1. 기존 허가 품목도 즉시 개정 기준을 적용해야 하나요?",
        "A1. 아닙니다. 시행일(2026.10.1) 당시 이미 허가된 품목은",
        "2027년 3월 31일까지 유예기간이 부여되며, 그 기간 동안은",
        "종전 기준 또는 개정 기준 중 하나를 선택하여 적용할 수 있습니다.",
        "유예기간 종료 이후에는 개정 기준을 반드시 적용해야 합니다.",
        "",
        "Q2. 유예기간 중 재시험을 하면 어떤 기준을 적용하나요?",
        "A2. 유예기간 중에는 품목별로 종전 기준 또는 개정 기준 중",
        "하나를 선택할 수 있으나, 동일 품목 내에서는 일관된 기준을",
        "적용해야 하며 임의로 혼용할 수 없습니다.",
        "",
        "Q3. 수출용 의약품에도 이 고시가 적용되나요?",
        "A3. 원칙적으로 국내 허가 품목에 한하여 적용됩니다.",
        "수출 전용 품목은 수입국 규정을 우선 적용하되,",
        "국내 GMP 실사 대상인 경우 자세한 사항은 원문 확인이 필요합니다.",
    ]
    text_obj = c.beginText(50, 800)
    text_obj.setFont(KOREAN_FONT, 10)
    for line in lines:
        text_obj.textLine(line)
    c.drawText(text_obj)
    c.save()
    return path


def make_seminar_pptx():
    path = os.path.join(SAMPLE_DIR, "seminar_sample.pptx")
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "2026년 7월 사내세미나: 미생물한도 개정 대응 방안"
    slide.placeholders[1].text_frame.text = "발표자: 김규제 책임연구원 / RA팀"

    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "개정 배경 및 유예기간 요약"
    tf = slide2.placeholders[1].text_frame
    tf.text = "총호기성생균수 기준 100->10 CFU/g 강화"
    p = tf.add_paragraph()
    p.text = "기존 허가품목 유예기간: 2027.3.31까지"
    p2 = tf.add_paragraph()
    p2.text = "QC팀 시험 배지 교체 준비 필요 (TSA 배지로 통일)"

    slide3 = prs.slides.add_slide(slide_layout)
    slide3.shapes.title.text = "RA팀 대응 체크리스트"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "허가변경 신청 대상 품목 리스트업 (8월까지)"
    p3 = tf3.add_paragraph()
    p3.text = "유예기간 내 시험성적서 기준 명시 필수"

    prs.save(path)
    return path


def main():
    db.init_db()

    notice_path = make_notice_pdf()
    notice_text = extract.extract_text(notice_path)
    db.insert_document(
        category="notice",
        title="완제의약품 미생물한도 시험법 개정 고시 (2026-45호, 예시)",
        extracted_text=notice_text,
        team_tag=None,
        uploader="관리자(샘플)",
        source_url="https://www.mfds.go.kr/",
        original_filename="notice_sample.pdf",
        stored_path=os.path.basename(notice_path),
    )
    # 실제 업로드 흐름과 동일하게 uploads 폴더로도 복사
    import shutil
    shutil.copy(notice_path, os.path.join(db.UPLOAD_DIR, os.path.basename(notice_path)))

    qna_path = make_qna_pdf()
    qna_text = extract.extract_text(qna_path)
    db.insert_document(
        category="qna",
        title="미생물한도 시험법 개정 관련 Q&A (예시)",
        extracted_text=qna_text,
        uploader="관리자(샘플)",
        source_url="https://www.mfds.go.kr/",
        original_filename="qna_sample.pdf",
        stored_path=os.path.basename(qna_path),
    )
    shutil.copy(qna_path, os.path.join(db.UPLOAD_DIR, os.path.basename(qna_path)))

    seminar_path = make_seminar_pptx()
    seminar_text = extract.extract_text(seminar_path)
    db.insert_document(
        category="seminar",
        title="2026년 7월 사내세미나: 미생물한도 개정 대응 방안",
        extracted_text=seminar_text,
        presenter="김규제 책임연구원",
        event_date="2026-07-05",
        uploader="관리자(샘플)",
        original_filename="seminar_sample.pptx",
        stored_path=os.path.basename(seminar_path),
    )
    shutil.copy(seminar_path, os.path.join(db.UPLOAD_DIR, os.path.basename(seminar_path)))

    print("샘플 데이터 3건 등록 완료.")


if __name__ == "__main__":
    main()
