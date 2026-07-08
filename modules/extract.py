"""
업로드된 파일(PDF/PPTX/TXT)에서 텍스트를 추출하는 모듈.
- 검색과 AI 분석은 여기서 추출된 텍스트를 기반으로 동작합니다.
"""
import os
from pypdf import PdfReader
from pptx import Presentation


def extract_text(file_path):
    """확장자에 따라 알맞은 추출기를 호출합니다."""
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".pdf":
            return _extract_pdf(file_path)
        elif ext in (".pptx", ".ppt"):
            return _extract_pptx(file_path)
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        else:
            return ""
    except Exception as e:
        return f"[텍스트 추출 실패: {e}]"


def _extract_pdf(file_path):
    reader = PdfReader(file_path)
    texts = []
    for page in reader.pages:
        try:
            texts.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n".join(texts).strip()


def _extract_pptx(file_path):
    prs = Presentation(file_path)
    texts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        slide_texts.append(line.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            slide_texts.append(cell.text.strip())
        if slide_texts:
            texts.append(f"[슬라이드 {i}] " + " / ".join(slide_texts))
    return "\n".join(texts).strip()
